import streamlit as st
import pdfplumber
import openai
import json
import pandas as pd
import io
import datetime

# 워드/PPT 라이브러리 불러오기 (에러 방지)
try:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    HAS_LIBS = True
except ImportError:
    HAS_LIBS = False

# ==========================================
# [설정] 본부장님의 API 키를 아래 따옴표 안에 넣어주세요.
API_KEY = st.secrets["OPENAI_API_KEY"]
# ==========================================

# 1. AI 분석 함수
def analyze_docs_with_gpt(text_data):
    if not API_KEY or "sk-" not in API_KEY:
        return {"error": "API 키가 설정되지 않았습니다."}
        
    client = openai.OpenAI(api_key=API_KEY)
    
    # 프롬프트 (안전하게 한 줄로 작성)
    prompt = "제공된 텍스트에서 재무 데이터(당기순이익, 전기순이익, 전전기순이익, 자산총계, 부채총계, 발행주식총수, 자본금)를 JSON으로 추출해. 단위는 원화 정수. 항목명: n_profit, n1_profit, n2_profit, total_assets, total_liabilities, total_shares, total_capital.\n\n" + text_data[:15000]

    try:
        response = client.chat.completions.create(
            model="gpt-4o", 
            messages=[
                {"role": "system", "content": "JSON으로만 답해."},
                {"role": "user", "content": prompt}
            ],
            response_format={"type": "json_object"}
        )
        return json.loads(response.choices[0].message.content)
    except Exception as e:
        return {"error": str(e)}

# 2. 세금 계산 함수
def calc_global_income_tax(income):
    if income <= 0: return 0
    tax = 0
    # 2024년 귀속 종합소득세율
    if income <= 14000000: tax = income * 0.06
    elif income <= 50000000: tax = (income * 0.15) - 1260000
    elif income <= 88000000: tax = (income * 0.24) - 5760000
    elif income <= 150000000: tax = (income * 0.35) - 15440000
    elif income <= 300000000: tax = (income * 0.38) - 19940000
    elif income <= 500000000: tax = (income * 0.40) - 25940000
    elif income <= 1000000000: tax = (income * 0.42) - 35940000
    else: tax = (income * 0.45) - 65440000
    return tax * 1.1

# 3. 전체 계산 로직
def calculate_all(n, n1, n2, assets, liab, shares, capital):
    if shares == 0: return None

    net_assets = assets - liab
    val_asset = max(0, net_assets / shares)
    
    w_profit = (n*3 + n1*2 + n2*1) / 6
    val_profit = max(0, (w_profit / shares) / 0.10)
    
    stock_value = (val_profit * 3 + val_asset * 2) / 5
    total_value = stock_value * shares
    face_value = capital / shares if shares > 0 else 5000 

    def calc_gift(amt):
        if amt <= 1e8: return amt * 0.10
        elif amt <= 5e8: return 1e7 + (amt - 1e8)*0.20
        elif amt <= 10e8: return 9e7 + (amt - 5e8)*0.30
        elif amt <= 30e8: return 2.4e8 + (amt - 10e8)*0.40
        else: return 10.4e8 + (amt - 30e8)*0.50
    
    tax_inherit = calc_gift(total_value)
    market_premium = max(0, stock_value - face_value)
    tax_sale = (market_premium * shares) * 0.20
    liquidation_income = max(0, net_assets - capital)
    tax_liquidation = calc_global_income_tax(liquidation_income)

    return {
        "val_asset": int(val_asset),
        "val_profit": int(val_profit),
        "stock_value": int(stock_value),
        "total_value": int(total_value),
        "tax_inherit": int(tax_inherit),
        "tax_sale": int(tax_sale),
        "liquidation_income": int(liquidation_income),
        "tax_liquidation": int(tax_liquidation)
    }

# 4. 워드 파일 생성 함수
def create_docx(company_name, res):
    doc = Document()
    title = f"{company_name} 기업가치 평가 보고서" if company_name else "기업가치 평가 보고서"
    heading = doc.add_heading(title, 0)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(f"작성일: {datetime.date.today()}").alignment = WD_ALIGN_PARAGRAPH.RIGHT
    doc.add_paragraph()

    doc.add_heading('1. 종합 평가 결과', level=1)
    p = doc.add_paragraph()
    runner = p.add_run(f"총 기업가치: {res['total_value']:,} 원")
    runner.bold = True
    doc.add_paragraph(f"1주당 평가액: {res['stock_value']:,} 원")
    doc.add_paragraph("(평가방법: 상증법상 보충적 평가 / 순손익 3 : 순자산 2)")

    doc.add_heading('2. 예상 세금 시뮬레이션', level=1)
    doc.add_paragraph(f"① 상속/증여 시: {res['tax_inherit']:,} 원")
    doc.add_paragraph(f"② 매각(양도) 시: {res['tax_sale']:,} 원 (세율 20%)")
    doc.add_paragraph(f"③ 청산 시 (종소세): {res['tax_liquidation']:,} 원")

    doc.add_heading('3. 상세 평가 내역', level=1)
    table = doc.add_table(rows=3, cols=2)
    table.style = 'Table Grid'
    table.rows[0].cells[0].text = "주당 순손익가치"
    table.rows[0].cells[1].text = f"{res['val_profit']:,} 원"
    table.rows[1].cells[0].text = "주당 순자산가치"
    table.rows[1].cells[1].text = f"{res['val_asset']:,} 원"
    table.rows[2].cells[0].text = "최종 평가액"
    table.rows[2].cells[1].text = f"{res['stock_value']:,} 원"

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# 5. PPT 파일 생성 함수 (슬라이드 양식)
def create_ppt(company_name, res):
    prs = Presentation()
    
    # [슬라이드 1] 표지
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"{company_name}\n기업가치 평가 보고서" if company_name else "기업가치 평가 보고서"
    subtitle.text = f"작성일: {datetime.date.today()}\n작성자: 기업 컨설팅 본부"

    # [슬라이드 2] 종합 평가 결과
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "1. 종합 평가 결과"
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    
    p = tf.add_paragraph()
    p.text = f"총 기업가치: {res['total_value']:,} 원"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(0, 51, 102) # 남색
    
    p = tf.add_paragraph()
    p.text = f"1주당 평가액: {res['stock_value']:,} 원"
    p.font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "평가 방식: 상속증여세법상 비상장주식 보충적 평가방법"
    p.level = 1

    # [슬라이드 3] 세금 시뮬레이션
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "2. 예상 세금 시뮬레이션"
    tf = slide.shapes.placeholders[1].text_frame
    
    p = tf.add_paragraph()
    p.text = "① 상속 / 증여 시 (예상)"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = f"▶ {res['tax_inherit']:,} 원 (공제 전 최대치)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "② 지분 매각 (양도) 시"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = f"▶ {res['tax_sale']:,} 원 (양도세율 20% 가정)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "③ 법인 청산 시 (종합소득세)"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = f"▶ {res['tax_liquidation']:,} 원"
    p.level = 1
    p = tf.add_paragraph()
    p.text = f"(배당소득 {res['liquidation_income']:,}원 기준)"
    p.level = 2

    # 파일 저장
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# 6. 웹 화면 구성
st.set_page_config(page_title="기업보고서", layout="wide")
st.title("📑 비상장주식 가치평가 시스템")

with st.sidebar:
    st.header("설정 및 입력")
    uploaded_files = st.file_uploader("PDF 파일 업로드", type=["pdf"], accept_multiple_files=True)
    if uploaded_files and st.button("문서 분석 실행"):
        with st.spinner("분석 중..."):
            full_text = ""
            for file in uploaded_files:
                with pdfplumber.open(file) as pdf:
                    for page in pdf.pages: full_text += page.extract_text() + "\n"
            data = analyze_docs_with_gpt(full_text)
            if "error" in data: st.error(data['error'])
            else:
                st.session_state['data'] = data
                st.success("완료")

company_name = st.text_input("회사명", placeholder="예: 삼성전자")

if 'data' not in st.session_state:
    st.info("왼쪽 사이드바에서 PDF를 업로드해주세요.")
else:
    with st.expander("데이터 확인 및 수정", expanded=True):
        data = st.session_state['data']
        c1, c2, c3 = st.columns(3)
        with c1:
            st.subheader("손익")
            n = st.number_input("당기순이익", value=int(data.get('n_profit', 0)))
            n1 = st.number_input("전기순이익", value=int(data.get('n1_profit', 0)))
            n2 = st.number_input("전전기순이익", value=int(data.get('n2_profit', 0)))
        with c2:
            st.subheader("재무")
            assets = st.number_input("자산총계", value=int(data.get('total_assets', 0)))
            liab = st.number_input("부채총계", value=int(data.get('total_liabilities', 0)))
        with c3:
            st.subheader("주식")
            shares = st.number_input("발행주식수", value=int(data.get('total_shares', 0)))
            capital = st.number_input("자본금", value=int(data.get('total_capital', 0)))

    st.divider()

    if st.button("📊 보고서 생성하기", use_container_width=True):
        res = calculate_all(n, n1, n2, assets, liab, shares, capital)
        
        if res:
            title_text = f"{company_name} 기업현황보고서" if company_name else "기업현황보고서"
            
            # HTML 보고서 디자인 (안전형)
            html_code = f"""
            <div style="background-color: white; padding: 20px; border-radius: 10px; border: 1px solid #ddd; box-shadow: 0 4px 6px rgba(0,0,0,0.1);">
                <h2 style="text-align: center; color: #2c3e50; border-bottom: 2px solid #2c3e50; padding-bottom: 10px;">{title_text}</h2>
                <p style="text-align: right; color: gray;">작성일: {datetime.date.today()}</p>
                
                <div style="background-color: #f0f7fb; padding: 20px; border-radius: 10px; border-left: 5px solid #3498db; margin-bottom: 20px;">
                    <h3 style="margin: 0; color: #2980b9;">🏢 총 기업가치: {res['total_value']:,} 원</h3>
                    <p style="margin: 5px 0 0 0; font-size: 1.1em;">1주당 평가액: <b>{res['stock_value']:,} 원</b></p>
                </div>

                <div style="display: flex; gap: 20px;">
                    <div style="flex: 1; background-color: #fff3cd; padding: 15px; border-radius: 10px; border: 1px solid #ffeeba;">
                        <h4 style="margin-top: 0; color: #856404;">💸 예상 세금 시뮬레이션</h4>
                        <p><b>1. 상속/증여 시:</b> <span style="color:#c0392b;">{res['tax_inherit']:,} 원</span></p>
                        <p><b>2. 매각(양도) 시:</b> <span style="color:#d35400;">{res['tax_sale']:,} 원</span></p>
                        <p><b>3. 청산(종소세):</b> <span style="color:#27ae60;">{res['tax_liquidation']:,} 원</span></p>
                    </div>
                    
                    <div style="flex: 1; background-color: #f8f9fa; padding: 15px; border-radius: 10px; border: 1px solid #ddd;">
                        <h4 style="margin-top: 0; color: #343a40;">📊 가치평가 상세</h4>
                        <p><b>주당 순손익가치:</b> {res['val_profit']:,} 원</p>
                        <p><b>주당 순자산가치:</b> {res['val_asset']:,} 원</p>
                        <p><b>최종 평가액:</b> {res['stock_value']:,} 원</p>
                    </div>
                </div>
            </div>
            """
            st.markdown(html_code, unsafe_allow_html=True)
            st.write("")

            # 다운로드 버튼 영역 (2단 구성)
            if HAS_LIBS:
                c1, c2 = st.columns(2)
                with c1:
                    docx = create_docx(company_name, res)
                    st.download_button("💾 워드(.docx) 다운로드", docx, f"{company_name}_보고서.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", type="primary", use_container_width=True)
                with c2:
                    ppt = create_ppt(company_name, res)
                    st.download_button("📺 PPT(.pptx) 다운로드", ppt, f"{company_name}_보고서.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", type="primary", use_container_width=True)
                
                st.info("💡 **팁:** 이미지(JPG)가 필요하시면, PPT 파일을 다운받아 열어서 **[파일] > [내보내기] > [파일 형식 변경] > [JPEG]**로 저장하시면 깔끔한 슬라이드 그림을 얻을 수 있습니다!")
            else:
                st.error("기능을 사용하려면 'pip install python-docx python-pptx' 설치가 필요합니다.")